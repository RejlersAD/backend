"""
Export Service (MODULE 12) — CSV / JSON / Primavera-ready CSV / Excel / PowerPoint.

CSV & JSON are the MVP-required minimum bar per spec. Excel, a
Primavera-compatible column layout, and a client-ready PowerPoint deck are
provided as they only need libraries already in requirements.txt (openpyxl,
python-pptx) — no extra dependency risk beyond python-pptx itself.
"""
from __future__ import annotations

import base64
import csv
import datetime
import io
import json
import os
import re
import uuid

from django.utils import timezone

from ..config import DEFAULT_CALENDAR, PRIMAVERA_XER_VERSION
from .calendar_utils import add_working_days

ACTIVITY_COLUMNS = [
    'id', 'wbs_code', 'name', 'discipline', 'deliverable', 'responsible_role',
    'workflow_status', 'original_duration_days', 'start_date', 'finish_date', 
    'total_float_days', 'is_critical', 'is_milestone',
]

PRIMAVERA_COLUMNS = [
    ('id', 'Activity ID'),
    ('name', 'Activity Name'),
    ('original_duration_days', 'Original Duration'),
    ('start_date', 'Start'),
    ('finish_date', 'Finish'),
    ('total_float_days', 'Total Float'),
    ('predecessors_str', 'Predecessors'),
]

# ─────────────────────────────────────────────────────────────────────────
# PowerPoint export — uses Rejlers' corporate template when the optional asset
# is available. A complete branded widescreen deck is generated from standard
# Office layouts when it is absent, so local and production exports remain
# functional without a binary template deployment step.
# ─────────────────────────────────────────────────────────────────────────
PPTX_TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), '..', 'assets', 'rejlers_template.pptx')

# Layout names inside rejlers_template.pptx — soft-coded so a template
# revision only needs a one-line update here (not a rewrite of the slide
# building logic) if Rejlers ever renames a layout in the master.
PPTX_LAYOUT_COVER = 'Start slide'
PPTX_LAYOUT_CONTENT = 'Text slide (light grey)'
PPTX_LAYOUT_TABLE = 'Graph/table slide (light grey)'
PPTX_LAYOUT_CLOSING = 'Last slide_Thanks'

# Placeholder indexes used by the above layouts (verified against the
# template — see idx values on `layout.placeholders`).
PPTX_HEADING_PLACEHOLDER_IDX = 12   # present on every content layout — slide heading
PPTX_BODY_PLACEHOLDER_IDX = 14      # "Text slide (light grey)" — bullet/body text area
PPTX_CONTENT_PLACEHOLDER_IDX = 15   # "Graph/table slide (light grey)" — table/graphic area

PPTX_MAX_TABLE_ROWS = 12  # keep slide tables readable — soft-coded cap



def _predecessors_str(activity: dict) -> str:
    return ', '.join(p['id'] for p in activity.get('predecessors', []) if p.get('id'))


def activities_to_csv(activities: list) -> str:
    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=ACTIVITY_COLUMNS, extrasaction='ignore')
    writer.writeheader()
    for a in activities:
        writer.writerow(a)
    return buf.getvalue()


def activities_to_primavera_csv(activities: list) -> str:
    buf = io.StringIO()
    headers = [label for _, label in PRIMAVERA_COLUMNS]
    writer = csv.writer(buf)
    writer.writerow(headers)
    for a in activities:
        row = dict(a)
        row['predecessors_str'] = _predecessors_str(a)
        writer.writerow([row.get(key, '') for key, _ in PRIMAVERA_COLUMNS])
    return buf.getvalue()


def eddr_to_csv(eddr: list) -> str:
    if not eddr:
        return ''
    fieldnames = sorted({key for row in eddr for key in row.keys()})
    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=fieldnames)
    writer.writeheader()
    writer.writerows(eddr)
    return buf.getvalue()


def generation_to_json(generation) -> str:
    return json.dumps({
        'project': generation.project.name,
        'version': generation.version,
        'intelligence': generation.intelligence,
        'wbs': generation.wbs,
        'activities': generation.activities,
        'logic_matrix': generation.logic_matrix,
        'eddr': generation.eddr,
        'milestones': generation.milestones,
        'manhours': generation.manhours,
        'validation': generation.validation,
        'narrative': generation.narrative,
    }, indent=2, default=str)


def activities_to_excel_bytes(activities: list) -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Activities'
    ws.append(ACTIVITY_COLUMNS)
    for a in activities:
        ws.append([a.get(col, '') for col in ACTIVITY_COLUMNS])

    stream = io.BytesIO()
    wb.save(stream)
    return stream.getvalue()


def _pptx_get_layout(prs, name):
    """Return a named corporate layout or a compatible standard layout.

    The corporate template is optional in deployments.  Falling back every
    slide to layout 0 (Title Slide) made content/table exports fail because
    those slides do not expose the expected placeholders.  Standard
    ``python-pptx`` layout indexes are stable for a newly-created deck.
    """
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    fallback_indexes = {
        PPTX_LAYOUT_COVER: 0,    # Title Slide
        PPTX_LAYOUT_CONTENT: 1,  # Title and Content
        PPTX_LAYOUT_TABLE: 5,    # Title Only
        PPTX_LAYOUT_CLOSING: 0,  # Title Slide
    }
    index = min(fallback_indexes.get(name, 5), len(prs.slide_layouts) - 1)
    return prs.slide_layouts[index]


def _pptx_placeholder(slide, idx):
    """Safely resolve a template placeholder by its stable idx."""
    try:
        return slide.placeholders[idx]
    except (KeyError, IndexError):
        return None


def _pptx_set_heading(slide, prs, heading, *, light=False):
    """Set a heading on either the corporate or default Office layout."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    shape = _pptx_placeholder(slide, PPTX_HEADING_PLACEHOLDER_IDX) or slide.shapes.title
    if shape is None:
        shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), prs.slide_width - Inches(1.6), Inches(0.8))
    shape.text = str(heading)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255) if light else RGBColor(15, 23, 42)
    return shape


def _pptx_body_placeholder(slide, prs):
    """Return a writable body shape, creating one when the layout has none."""
    from pptx.util import Inches

    body = _pptx_placeholder(slide, PPTX_BODY_PLACEHOLDER_IDX)
    if body is None:
        title = slide.shapes.title
        title_shape_id = getattr(title, 'shape_id', None)
        body = next(
            (
                shape for shape in slide.placeholders
                if getattr(shape, 'shape_id', None) != title_shape_id
                and getattr(shape, 'has_text_frame', False)
            ),
            None,
        )
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), prs.slide_width - Inches(1.8), prs.slide_height - Inches(2.15))
    return body


def _pptx_add_footer(slide, prs, text='RADAI Project Planning Application'):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    footer = slide.shapes.add_textbox(Inches(0.8), prs.slide_height - Inches(0.42), prs.slide_width - Inches(1.6), Inches(0.22))
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = 2
    for run in paragraph.runs:
        run.font.name = 'Arial'
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(100, 116, 139)


def _pptx_strip_sample_slides(prs):
    """rejlers_template.pptx ships with sample/demo slides (used internally by
    Rejlers to showcase the template). Remove them so generated decks start
    clean while keeping every master/layout/theme/logo asset intact."""
    from pptx.oxml.ns import qn

    slide_id_list = prs.slides._sldIdLst
    for sld in list(slide_id_list):
        r_id = sld.get(qn('r:id'))
        prs.part.drop_rel(r_id)
        slide_id_list.remove(sld)


def _pptx_add_cover_slide(prs, project, generation):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_pptx_get_layout(prs, PPTX_LAYOUT_COVER))
    corporate_heading = _pptx_placeholder(slide, PPTX_HEADING_PLACEHOLDER_IDX)
    if corporate_heading is not None:
        tf = corporate_heading.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = project.name or 'Untitled Planning Project'
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        heading = _pptx_set_heading(slide, prs, project.name or 'Untitled Planning Project', light=True)
        heading.left = Inches(0.9)
        heading.top = Inches(1.4)
        heading.width = prs.slide_width - Inches(1.8)
        heading.height = Inches(1.25)
        tf = _pptx_body_placeholder(slide, prs).text_frame
        tf.clear()
        tf.word_wrap = True

    subtitle_bits = [b for b in [project.phase, project.client, project.location] if b]
    lines = (
        ' · '.join(subtitle_bits) if subtitle_bits else 'Project Planning Presentation',
        f'Effective Date: {project.effective_date or "—"}',
        f'RADAI Project Planning Application · Schedule v{generation.version}',
    )
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 and corporate_heading is None else tf.add_paragraph()
        p.text = line
        if corporate_heading is None:
            for run in p.runs:
                run.font.name = 'Arial'
                run.font.size = Pt(18 if i == 0 else 13)
                run.font.color.rgb = RGBColor(203, 213, 225)
    return slide


def _pptx_add_bullet_slide(prs, heading, bullets, max_items=10):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    slide = prs.slides.add_slide(_pptx_get_layout(prs, PPTX_LAYOUT_CONTENT))
    _pptx_set_heading(slide, prs, heading)

    tf = _pptx_body_placeholder(slide, prs).text_frame
    tf.clear()
    tf.word_wrap = True
    items = bullets[:max_items] or ['No data available.']
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(text)
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(51, 65, 85)
    if len(bullets) > max_items:
        p = tf.add_paragraph()
        p.text = f'…and {len(bullets) - max_items} more (see full export for details).'
    _pptx_add_footer(slide, prs)
    return slide


def _pptx_paginate_paragraphs(paragraphs, max_chars=750, max_items=5):
    """Split narrative text into slide-sized pages without dropping content.

    PowerPoint does not reliably auto-create continuation slides when text
    overflows a placeholder.  Keep a conservative character and paragraph
    budget, splitting unusually long paragraphs at sentence/word boundaries.
    """
    segments = []
    for paragraph in paragraphs:
        text = re.sub(r'\s+', ' ', str(paragraph)).strip()
        if not text:
            continue

        sentences = [part.strip() for part in re.split(r'(?<=[.!?])\s+', text) if part.strip()]
        current = ''
        for sentence in sentences or [text]:
            if len(sentence) > max_chars:
                words = sentence.split()
                word_chunk = ''
                for word in words:
                    candidate = f'{word_chunk} {word}'.strip()
                    if word_chunk and len(candidate) > max_chars:
                        if current:
                            segments.append(current)
                            current = ''
                        segments.append(word_chunk)
                        word_chunk = word
                    else:
                        word_chunk = candidate
                sentence_parts = [word_chunk] if word_chunk else []
            else:
                sentence_parts = [sentence]

            for part in sentence_parts:
                candidate = f'{current} {part}'.strip()
                if current and len(candidate) > max_chars:
                    segments.append(current)
                    current = part
                else:
                    current = candidate
        if current:
            segments.append(current)

    pages = []
    page = []
    page_chars = 0
    for segment in segments:
        required = len(segment) + (1 if page else 0)
        if page and (len(page) >= max_items or page_chars + required > max_chars):
            pages.append(page)
            page = []
            page_chars = 0
        page.append(segment)
        page_chars += len(segment) + (1 if page_chars else 0)
    if page:
        pages.append(page)
    return pages


def _pptx_add_table_slide(prs, heading, headers, rows, max_rows=PPTX_MAX_TABLE_ROWS, footnote=None):
    """Adds a table-content slide using the template's 'Graph/table slide'
    layout. The layout's empty content placeholder only supports simple text
    (it is not a native table placeholder in python-pptx), so we read its
    on-brand position/size then swap it for a real table at that exact spot —
    keeping the table perfectly aligned with the template's grid."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt, Emu, Inches

    slide = prs.slides.add_slide(_pptx_get_layout(prs, PPTX_LAYOUT_TABLE))
    _pptx_set_heading(slide, prs, heading)

    content_ph = _pptx_placeholder(slide, PPTX_CONTENT_PLACEHOLDER_IDX)
    if content_ph is not None:
        left, top, width, height = content_ph.left, content_ph.top, content_ph.width, content_ph.height
        content_ph._element.getparent().remove(content_ph._element)
    else:
        left, top = Inches(0.75), Inches(1.45)
        width, height = prs.slide_width - Inches(1.5), prs.slide_height - Inches(2.05)

    note_lines = []
    if len(rows) > max_rows:
        note_lines.append(f'Showing {max_rows} of {len(rows)} rows — full detail available via CSV/Excel export.')
    if footnote:
        note_lines.append(footnote)
    table_height = Emu(int(height * 0.85)) if note_lines else height

    display_rows = rows[:max_rows]
    graphic_frame = slide.shapes.add_table(len(display_rows) + 1, len(headers), left, top, width, table_height)
    table = graphic_frame.table
    for c, header in enumerate(headers):
        table.cell(0, c).text = str(header)
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = RGBColor(30, 64, 175)
        for run in table.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = Pt(11)
    for r, row in enumerate(display_rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = '' if value is None else str(value)
            if r % 2 == 0:
                table.cell(r, c).fill.solid()
                table.cell(r, c).fill.fore_color.rgb = RGBColor(241, 245, 249)
            for run in table.cell(r, c).text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(51, 65, 85)

    if note_lines:
        note_top = Emu(int(top) + int(table_height) + Emu(int(prs.slide_height * 0.015)))
        note_box = slide.shapes.add_textbox(left, note_top, width, Emu(int(prs.slide_height * 0.08)))
        ntf = note_box.text_frame
        ntf.word_wrap = True
        for i, line in enumerate(note_lines):
            p = ntf.paragraphs[0] if i == 0 else ntf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(11)
            run.font.italic = True
    _pptx_add_footer(slide, prs)
    return slide


def _pptx_add_closing_slide(prs, message='Thank You'):
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(_pptx_get_layout(prs, PPTX_LAYOUT_CLOSING))
    if _pptx_placeholder(slide, PPTX_HEADING_PLACEHOLDER_IDX) is not None:
        _pptx_set_heading(slide, prs, message)
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        _pptx_set_heading(slide, prs, message, light=True)
    return slide


def generation_to_pptx_bytes(generation) -> bytes:
    """
    Builds a client/internal-review-ready PowerPoint deck summarizing a
    PlanningGeneration. Rejlers' corporate template is used when available;
    otherwise safe standard layouts produce a complete branded widescreen
    deck with cover, agenda, content, table and closing slides.
    """
    from pptx import Presentation
    from pptx.util import Inches

    project = generation.project
    activities = generation.activities or []
    wbs = generation.wbs or []
    eddr = generation.eddr or []
    manhours = generation.manhours or {}
    validation = generation.validation or []
    intelligence = generation.intelligence or {}

    template_loaded = False
    if os.path.isfile(PPTX_TEMPLATE_PATH):
        try:
            prs = Presentation(PPTX_TEMPLATE_PATH)
            _pptx_strip_sample_slides(prs)
            template_loaded = True
        except Exception:
            prs = Presentation()
    else:
        prs = Presentation()

    # A generated deck must remain functional even when no corporate template
    # asset is shipped. The blank-deck path uses a modern 16:9 canvas and the
    # helper functions above supply safe layouts, text boxes and branding.
    if not template_loaded:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # 1) Cover
    _pptx_add_cover_slide(prs, project, generation)

    # 2) Agenda
    milestones = generation.milestones or [a for a in activities if a.get('is_milestone')]
    agenda_items = ['Project Overview', 'Work Breakdown Structure', 'Schedule Summary']
    if milestones:
        agenda_items.append('Key Milestones')
    agenda_items += [
        'Engineering Document Deliverable Register', 'Manhour Estimate',
        'Validation & Quality Checks', 'Executive Summary',
    ]
    _pptx_add_bullet_slide(prs, 'Agenda', agenda_items, max_items=len(agenda_items))

    # 3) Project Overview
    overview_bullets = [
        f'Client: {project.client or "—"}',
        f'Location: {project.location or "—"}',
        f'Phase: {project.phase or "—"}',
        f'Effective Date: {project.effective_date or "—"}',
        f'Planned Duration: {project.duration_months} month(s)',
        f'Reference Documents Analyzed: {len(intelligence.get("disciplines", {}) or {})} discipline(s)',
        f'Schedule Version: v{generation.version}',
    ]
    _pptx_add_bullet_slide(prs, 'Project Overview', overview_bullets)

    # 4) WBS Summary
    top_level = [n for n in wbs if n.get('level') == 1] or wbs
    wbs_bullets = [f'{n.get("code", "")} — {n.get("name", "")}' for n in top_level]
    if not wbs_bullets:
        wbs_bullets = ['WBS not yet generated.']
    else:
        wbs_bullets.insert(0, f'{len(wbs)} total WBS node(s) across the project.')
    _pptx_add_bullet_slide(prs, 'Work Breakdown Structure', wbs_bullets)

    # 5) Schedule Summary
    critical = [a for a in activities if a.get('is_critical')]
    finish_dates = [a.get('finish_date') for a in activities if a.get('finish_date')]
    schedule_bullets = [
        f'Total Activities: {len(activities)}',
        f'Critical Path Activities: {len(critical)}',
        f'Milestones: {len(milestones)}',
        f'Projected Finish: {max(finish_dates) if finish_dates else "—"}',
    ]
    _pptx_add_bullet_slide(prs, 'Schedule Summary', schedule_bullets, max_items=6)

    # 6) Key Milestones (only when milestones exist)
    if milestones:
        _pptx_add_table_slide(
            prs, 'Key Milestones',
            ['Milestone', 'Finish Date'],
            [[m.get('name', ''), m.get('finish_date', '')] for m in milestones],
            max_rows=8,
        )

    # 7) EDDR Summary
    if eddr:
        _pptx_add_table_slide(
            prs, 'Engineering Document Deliverable Register',
            ['Discipline', 'Deliverable', 'Final Issue'],
            [[row.get('discipline', ''), row.get('deliverable_name', ''), row.get('final_issue_date', '')] for row in eddr],
        )
    else:
        _pptx_add_bullet_slide(prs, 'Engineering Document Deliverable Register', ['EDDR not yet generated.'])

    # 8) Manhours Summary
    by_discipline = manhours.get('by_discipline') or []
    if by_discipline:
        grand_total_hours = manhours.get('grand_total_man_hours')
        grand_total_days = sum(r.get('man_days') or 0 for r in by_discipline)
        footnote = (
            f'Grand Total: {grand_total_days} man-days / {grand_total_hours} man-hours'
            if grand_total_hours is not None else None
        )
        _pptx_add_table_slide(
            prs, 'Manhour Estimate',
            ['Discipline', 'Role', 'Man-Days', 'Man-Hours'],
            [[r.get('discipline_name', ''), r.get('responsible_role', ''), r.get('man_days', ''), r.get('man_hours', '')] for r in by_discipline],
            max_rows=8, footnote=footnote,
        )
    else:
        _pptx_add_bullet_slide(prs, 'Manhour Estimate', ['Manhour estimate not yet generated.'])

    # 9) Validation Summary
    counts = {}
    for issue in validation:
        counts[issue.get('severity', 'pass')] = counts.get(issue.get('severity', 'pass'), 0) + 1
    validation_bullets = [
        f'Pass: {counts.get("pass", 0)}',
        f'Warnings: {counts.get("warning", 0)}',
        f'Critical Issues: {counts.get("critical", 0)}',
    ]
    critical_msgs = [i.get('message', '') for i in validation if i.get('severity') == 'critical']
    if critical_msgs:
        validation_bullets.append('Key issues:')
        validation_bullets.extend(critical_msgs[:5])
    _pptx_add_bullet_slide(prs, 'Validation & Quality Checks', validation_bullets, max_items=10)

    # 10) Narrative / Executive Summary
    narrative_text = (generation.narrative or '').strip()
    paragraphs = [p.strip() for p in narrative_text.split('\n') if p.strip()]
    summary_pages = _pptx_paginate_paragraphs(
        paragraphs or ['Narrative not yet generated.'],
    )
    total_summary_pages = len(summary_pages)
    for page_number, summary_bullets in enumerate(summary_pages, start=1):
        heading = 'Executive Summary'
        if total_summary_pages > 1:
            heading = f'Executive Summary ({page_number}/{total_summary_pages})'
        _pptx_add_bullet_slide(
            prs,
            heading,
            summary_bullets,
            max_items=len(summary_bullets),
        )

    # 11) Closing
    _pptx_add_closing_slide(prs, 'Thank You')

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()


# ─────────────────────────────────────────────────────────────────────────
# Primavera P6 .xer export — schedule-only subset (CURRTYPE / PROJECT /
# CALENDAR / PROJWBS / TASK / TASKPRED). Field layout is soft-coded against
# the P6 XER schema used by the reference sample files supplied for this
# feature (Documents/Project Control/Planning Package/*.xer) — every %F/%R
# column list below mirrors that schema so the exported file opens cleanly
# in Primavera P6 (File > Import > Primavera PM - (XER)). The ERMHDR version
# string itself is soft-coded in ../config.py (PRIMAVERA_XER_VERSION) since
# it must match the target site's installed P6 version, not the schema layout.
# ─────────────────────────────────────────────────────────────────────────
XER_VERSION = PRIMAVERA_XER_VERSION
XER_PROJ_ID = 1
XER_CLNDR_ID = 1
XER_LINE_SEP = '\r\n'  # native XER line ending

# Primavera P6 enforces hard character limits on several short-identifier
# fields at the database schema level. Exceeding them doesn't just truncate
# on import — P6's XER importer silently REJECTS the offending row (and,
# critically, when it's the PROJECT row itself, the entire project — WBS,
# activities, everything — never gets created, so the import appears to
# do nothing at all). Every value written into these fields below is
# capped against these soft-coded limits so that can never happen again.
XER_PROJ_SHORT_NAME_MAX_LEN = 20   # PROJECT.proj_short_name ("Project ID")
XER_WBS_SHORT_NAME_MAX_LEN = 20    # PROJWBS.wbs_short_name ("WBS Code")
XER_WBS_NAME_MAX_LEN = 100         # PROJWBS.wbs_name
XER_TASK_CODE_MAX_LEN = 20         # TASK.task_code ("Activity ID")
XER_TASK_NAME_MAX_LEN = 120        # TASK.task_name ("Activity Name")

_XER_PRED_TYPE = {'FS': 'PR_FS', 'SS': 'PR_SS', 'FF': 'PR_FF', 'SF': 'PR_SF'}


def _derive_proj_short_name(project) -> str:
    """P6's Project ID is a short *code*, not the descriptive project name
    (the reference sample files use e.g. 'Q-101685') — and it's capped at
    `XER_PROJ_SHORT_NAME_MAX_LEN` characters. Prefer a trailing bracketed
    reference number if the project name has one (e.g. '... (30201-50198)'
    -> '30201-50198'), otherwise sanitize the name into a code-safe token.
    """
    name = (project.name or '').strip()
    match = re.search(r'\(([A-Za-z0-9][A-Za-z0-9\-_/]{2,})\)\s*$', name)
    candidate = match.group(1) if match else name
    # P6 Project IDs are conventionally alphanumeric + hyphen/underscore —
    # collapse anything else (spaces, punctuation) into a single hyphen.
    candidate = re.sub(r'[^A-Za-z0-9\-_]+', '-', candidate).strip('-')
    candidate = re.sub(r'-{2,}', '-', candidate)
    return (candidate or 'PROJECT')[:XER_PROJ_SHORT_NAME_MAX_LEN]


def _xer_field(value) -> str:
    """Sanitizes a single tab-delimited field value (XER is TSV-based, one
    record per line — tabs/newlines inside a value would corrupt the file)."""
    if value is None:
        return ''
    return str(value).replace('\t', ' ').replace('\r', ' ').replace('\n', ' ').strip()


def _xer_row(*values) -> str:
    return '\t'.join(['%R'] + [_xer_field(v) for v in values])


def _xer_datetime(date_str, time_str='08:00') -> str:
    if not date_str:
        return ''
    return f'{date_str} {time_str}'


def _xer_calendar_data(hours_per_day: int) -> str:
    """Builds a Mon–Fri working calendar (Sat/Sun off) at the given daily
    hour span, in P6's proprietary calendar syntax — mirrors the 'Standard
    5 Day Workweek' calendar found in the reference sample .xer files."""
    start_hr = '08:00'
    # P6 clock is 24h; workday finish = start + hours_per_day. Previous
    # code only handled hours_per_day < 16 and produced garbage on 8h days
    # (2×8=16 → '16:00' vs P6's expected end-of-shift), so compute directly.
    end_total_min = 8 * 60 + int(hours_per_day * 60)
    end_hr = f'{(end_total_min // 60) % 24:02d}:{end_total_min % 60:02d}'
    workday = f'(0||0(s|{start_hr}|f|{end_hr})())'
    days = ''.join(
        f'(0||{d}()({workday if 2 <= d <= 6 else ""}))'
        for d in range(1, 8)
    )
    return (
        f'(0||CalendarData()(  (0||DaysOfWeek()(    {days}))  '
        f'(0||VIEW(ShowTotal|Y)())  (0||Exceptions()())))'
    )


# ─────────────────────────────────────────────────────────────────────────
# XER — GUID + soft-coded structural identifiers
# P6 uses base64-encoded 16-byte UUIDs (22 chars, no '=' padding) as record
# GUIDs across TASK / PROJWBS / PROJECT. Generating real GUIDs (instead of
# leaving the column blank) keeps update semantics correct on re-import and
# avoids "duplicate" collisions with any prior imported schedule.
# ─────────────────────────────────────────────────────────────────────────
XER_OBS_ID = 1
XER_OBS_NAME = 'Enterprise'
XER_FINTMPL_ID = 1
XER_FINTMPL_NAME = 'Calendar'
XER_SCHEDOPTIONS_ID = 1
# Default working shift (soft-coded here — matches sample XER 08:00–17:00).
XER_WORK_START = '08:00'
XER_WORK_END = '17:00'


def _xer_guid() -> str:
    """Returns a Primavera-style 22-character base64 GUID (no '=' padding)."""
    return base64.b64encode(uuid.uuid4().bytes).decode('ascii').rstrip('=')


def _xer_add_working_days(iso_date: str, days: int, end_time: str = XER_WORK_END) -> str:
    """Shifts `iso_date` by `days` working days (Mon–Fri) and returns a
    P6 datetime literal `YYYY-MM-DD HH:MM`. Returns '' if the input date
    can't be parsed."""
    if not iso_date:
        return ''
    try:
        base = datetime.date.fromisoformat(iso_date)
    except (ValueError, TypeError):
        return ''
    try:
        shifted = add_working_days(base, int(days))
    except (ValueError, TypeError):
        return ''
    return f'{shifted.isoformat()} {end_time}'


def generation_to_xer_bytes(generation) -> bytes:
    """
    Exports this generation's WBS + Activities + Logic Matrix as a
    Primavera P6-compatible .xer schedule file.

    Layout mirrors a native P6-exported schedule (Documents/Project Control/
    Planning Package/Sample.txt): every mandatory section — ERMHDR, CURRTYPE,
    FINTMPL, OBS, PROJECT, CALENDAR, SCHEDOPTIONS, PROJWBS, TASK, TASKPRED —
    is emitted with the full column set P6 expects, so File > Import >
    Primavera PM (XER) succeeds and the imported project retains WBS
    hierarchy, activity dates, constraint types, GUIDs and full FS/SS/FF/SF
    predecessor logic.

    Returns bytes encoded as Windows-1252 (P6's native XER encoding) with
    UTF-8 as a defensive fallback for characters that don't map (e.g.
    non-Latin project names).
    """
    project = generation.project
    wbs_nodes = generation.wbs or []
    activities = generation.activities or []
    logic_matrix = generation.logic_matrix or []

    calendar = dict(DEFAULT_CALENDAR)
    calendar.update(getattr(project, 'calendar_overrides', None) or {})
    hours_per_day = int(calendar.get('hours_per_day') or 8)
    days_per_week = int(calendar.get('working_days_per_week') or 5)

    now = timezone.now()
    now_stamp = now.strftime('%Y-%m-%d %H:%M')
    export_date = now.strftime('%Y-%m-%d')
    export_user = (getattr(generation.generated_by, 'username', '') or 'radai').strip() or 'radai'
    export_user_full = getattr(generation.generated_by, 'get_full_name', lambda: '')() or export_user

    # WBS code -> synthetic numeric wbs_id (P6 requires numeric FKs).
    wbs_id_map = {node['code']: 1000 + i for i, node in enumerate(wbs_nodes, start=1)}
    wbs_guid_map = {node['code']: _xer_guid() for node in wbs_nodes}
    root_wbs_id = wbs_id_map.get('1') or (next(iter(wbs_id_map.values())) if wbs_id_map else 1000)

    # Activity id (e.g. "PR-100") -> synthetic numeric task_id.
    task_id_map = {a['id']: 2000 + i for i, a in enumerate(activities, start=1)}
    task_guid_map = {a['id']: _xer_guid() for a in activities}

    starts = [a['start_date'] for a in activities if a.get('start_date')]
    finishes = [a['finish_date'] for a in activities if a.get('finish_date')]
    plan_start = min(starts) if starts else (project.effective_date.isoformat() if project.effective_date else export_date)
    plan_end = max(finishes) if finishes else plan_start

    week_hours = hours_per_day * days_per_week
    month_hours = week_hours * 4        # P6 convention: 4 weeks/month
    year_hours = week_hours * 52

    lines = [
        f'ERMHDR\t{XER_VERSION}\t{export_date}\tProject\t{export_user}\t{export_user_full}\t'
        f'dbxDatabaseNoName\tProject Management\tUSD'
    ]

    # ── CURRTYPE ──
    lines.append('%T\tCURRTYPE')
    lines.append('%F\tcurr_id\tdecimal_digit_cnt\tcurr_symbol\tdecimal_symbol\tdigit_group_symbol\t'
                  'pos_curr_fmt_type\tneg_curr_fmt_type\tcurr_type\tcurr_short_name\tgroup_digit_cnt\tbase_exch_rate')
    lines.append(_xer_row(1, 2, '$', '.', ',', '#1.1', '(#1.1)', 'US Dollar', 'USD', 3, 1))

    # ── FINTMPL (Finance Template — required by P6 PROJECT.fintmpl_id FK) ──
    lines.append('%T\tFINTMPL')
    lines.append('%F\tfintmpl_id\tfintmpl_name\tdefault_flag')
    lines.append(_xer_row(XER_FINTMPL_ID, XER_FINTMPL_NAME, 'Y'))

    # ── OBS (Organizational Breakdown Structure — required by PROJWBS.obs_id) ──
    lines.append('%T\tOBS')
    lines.append('%F\tobs_id\tparent_obs_id\tguid\tseq_num\tobs_name\tobs_descr')
    lines.append(_xer_row(XER_OBS_ID, '', _xer_guid(), 0, XER_OBS_NAME, ''))

    # ── PROJECT ──
    # Column list mirrors Sample.txt exactly so every downstream P6
    # subsystem (baseline mgr, resource loader, scheduler) finds the
    # values it expects; blank strings are valid P6 "unset".
    lines.append('%T\tPROJECT')
    lines.append('%F\tproj_id\tfy_start_month_num\trsrc_self_add_flag\tallow_complete_flag\t'
                  'rsrc_multi_assign_flag\tcheckout_flag\tproject_flag\tstep_complete_flag\t'
                  'cost_qty_recalc_flag\tbatch_sum_flag\tname_sep_char\tdef_complete_pct_type\t'
                  'proj_short_name\tacct_id\torig_proj_id\tsource_proj_id\tbase_type_id\t'
                  'clndr_id\tsum_base_proj_id\ttask_code_base\ttask_code_step\tpriority_num\t'
                  'wbs_max_sum_level\tstrgy_priority_num\tlast_checksum\tcritical_drtn_hr_cnt\t'
                  'def_cost_per_qty\tlast_recalc_date\tplan_start_date\tplan_end_date\t'
                  'scd_end_date\tadd_date\tlast_tasksum_date\tfcst_start_date\tdef_duration_type\t'
                  'task_code_prefix\tguid\tdef_qty_type\tadd_by_name\tweb_local_root_path\t'
                  'proj_url\tdef_rate_type\tadd_act_remain_flag\tact_this_per_link_flag\t'
                  'def_task_type\tact_pct_link_flag\tcritical_path_type\ttask_code_prefix_flag\t'
                  'def_rollup_dates_flag\tuse_project_baseline_flag\trem_target_link_flag\t'
                  'reset_planned_flag\tallow_neg_act_flag\tsum_assign_level\tlast_fin_dates_id\t'
                  'fintmpl_id\tlast_baseline_update_date\tcr_external_key\tapply_actuals_date\t'
                  'location_id\tloaded_scope_level\texport_flag\tnew_fin_dates_id\t'
                  'baselines_to_export\tbaseline_names_to_export\tnext_data_date\t'
                  'close_period_flag\tsum_refresh_date\ttrsrcsum_loaded\tsumtask_loaded')
    lines.append(_xer_row(
        XER_PROJ_ID, 1, 'Y', 'Y', 'Y', 'N', 'Y', 'N', 'N', 'Y', '.', 'CP_Drtn',
        _derive_proj_short_name(project),
        '', '', '', '',                       # acct_id, orig_proj_id, source_proj_id, base_type_id
        XER_CLNDR_ID,
        '',                                   # sum_base_proj_id
        1000, 10, 10, 2, 500,                 # code_base, code_step, priority_num, wbs_max_sum_level, strgy_priority_num
        '', 0, '0.0000',                      # last_checksum, critical_drtn_hr_cnt, def_cost_per_qty
        _xer_datetime(plan_start, '00:00'),   # last_recalc_date
        _xer_datetime(plan_start, '00:00'),
        _xer_datetime(plan_end, '00:00'),
        _xer_datetime(plan_end, '00:00'),     # scd_end_date
        now_stamp, '', '',                    # add_date, last_tasksum_date, fcst_start_date
        'DT_FixedDUR2', 'A', _xer_guid(), 'QT_Hour', export_user, '', '',
        'COST_PER_QTY', 'N', 'Y', 'TT_Task', 'Y',
        'CT_TotFloat', 'Y', 'Y', 'Y', 'Y', 'N', 'N',
        'SL_Taskrsrc', '',                    # sum_assign_level, last_fin_dates_id
        XER_FINTMPL_ID, '', '', '', '',       # fintmpl_id, last_baseline_update_date, cr_external_key, apply_actuals_date, location_id
        7, 'Y', '', '', '', '', '', '', '', '',
    ))

    # ── CALENDAR ──
    lines.append('%T\tCALENDAR')
    lines.append('%F\tclndr_id\tdefault_flag\tclndr_name\tproj_id\tbase_clndr_id\tlast_chng_date\t'
                  'clndr_type\tday_hr_cnt\tweek_hr_cnt\tmonth_hr_cnt\tyear_hr_cnt\trsrc_private\tclndr_data')
    lines.append(_xer_row(
        # default_flag='N' — this calendar is project-scoped, not a global
        # default P6 database calendar.
        XER_CLNDR_ID, 'N', 'Standard 5 Day Workweek', '', '', now_stamp,
        'CA_Base', hours_per_day, week_hours, month_hours, year_hours,
        'N', _xer_calendar_data(hours_per_day),
    ))

    # ── SCHEDOPTIONS (required for CPM scheduling on import) ──
    lines.append('%T\tSCHEDOPTIONS')
    lines.append('%F\tschedoptions_id\tproj_id\tsched_outer_depend_type\tsched_open_critical_flag\t'
                  'sched_lag_early_start_flag\tsched_retained_logic\tsched_setplantoforecast\t'
                  'sched_float_type\tsched_calendar_on_relationship_lag\tsched_use_expect_end_flag\t'
                  'sched_progress_override\tlevel_float_thrs_cnt\tlevel_outer_assign_flag\t'
                  'level_outer_assign_priority\tlevel_over_alloc_pct\tlevel_within_float_flag\t'
                  'level_keep_sched_date_flag\tlevel_all_rsrc_flag\tsched_use_project_end_date_for_float\t'
                  'enable_multiple_longest_path_calc\tlimit_multiple_longest_path_calc\t'
                  'max_multiple_longest_path\tuse_total_float_multiple_longest_paths\t'
                  'key_activity_for_multiple_longest_paths\tLevelPriorityList')
    lines.append(_xer_row(
        XER_SCHEDOPTIONS_ID, XER_PROJ_ID,
        'SD_Both', 'N', 'Y', 'Y', 'N', 'FT_FF', 'rcal_Predecessor', 'Y', 'N',
        0, 'N', 5, 25, 'N', 'Y', 'Y', 'Y', 'N', 'Y', 10, 'Y', '',
        'priority_type,ASC_BY_FIELD/ASC',
    ))

    # ── PROJWBS ──
    lines.append('%T\tPROJWBS')
    lines.append('%F\twbs_id\tproj_id\tobs_id\tseq_num\test_wt\tproj_node_flag\tsum_data_flag\t'
                  'status_code\twbs_short_name\twbs_name\tphase_id\tparent_wbs_id\tev_user_pct\t'
                  'ev_etc_user_value\torig_cost\tindep_remain_total_cost\tann_dscnt_rate_pct\t'
                  'dscnt_period_type\tindep_remain_work_qty\tanticip_start_date\t'
                  'anticip_end_date\tev_compute_type\tev_etc_compute_type\tguid\ttmpl_guid\t'
                  'plan_open_state')
    for i, node in enumerate(wbs_nodes, start=1):
        parent_code = node.get('parent_code')
        lines.append(_xer_row(
            wbs_id_map[node['code']], XER_PROJ_ID, XER_OBS_ID, i, 1,
            'Y' if node.get('level') == 0 else 'N', 'N', 'WS_Open',
            node['code'].split('.')[-1][:XER_WBS_SHORT_NAME_MAX_LEN],
            (node.get('name') or '')[:XER_WBS_NAME_MAX_LEN],
            '',                               # phase_id
            wbs_id_map.get(parent_code, '') if parent_code else '',
            6, '0.88', '0.0000', '0.0000',   # EV defaults matching Sample.txt
            '', '', '', '', '', '',
            'EC_Cmp_pct', 'EE_Rem_hr',
            wbs_guid_map[node['code']], '', '',
        ))

    # ── TASK ──
    # Full column set from the sample XER — every date/flag/priority
    # field is emitted so P6 doesn't re-default them on import.
    lines.append('%T\tTASK')
    lines.append('%F\ttask_id\tproj_id\twbs_id\tclndr_id\tphys_complete_pct\trev_fdbk_flag\t'
                  'est_wt\tlock_plan_flag\tauto_compute_act_flag\tcomplete_pct_type\ttask_type\t'
                  'duration_type\tstatus_code\ttask_code\ttask_name\trsrc_id\ttotal_float_hr_cnt\t'
                  'free_float_hr_cnt\tremain_drtn_hr_cnt\tact_work_qty\tremain_work_qty\t'
                  'target_work_qty\ttarget_drtn_hr_cnt\ttarget_equip_qty\tact_equip_qty\t'
                  'remain_equip_qty\tcstr_date\tact_start_date\tact_end_date\tlate_start_date\t'
                  'late_end_date\texpect_end_date\tearly_start_date\tearly_end_date\t'
                  'restart_date\treend_date\ttarget_start_date\ttarget_end_date\t'
                  'rem_late_start_date\trem_late_end_date\tcstr_type\tpriority_type\t'
                  'suspend_date\tresume_date\tfloat_path\tfloat_path_order\tguid\ttmpl_guid\t'
                  'cstr_date2\tcstr_type2\tdriving_path_flag\tact_this_per_work_qty\t'
                  'act_this_per_equip_qty\texternal_early_start_date\texternal_late_end_date\t'
                  'create_date\tupdate_date\tcreate_user\tupdate_user\tlocation_id\tcrt_path_num')
    for activity in activities:
        wbs_id = wbs_id_map.get(activity.get('wbs_code'), root_wbs_id)
        is_milestone = bool(activity.get('is_milestone'))
        duration_hrs = int((activity.get('original_duration_days') or 0) * hours_per_day)
        # P6 milestones must always have zero duration; force it here so
        # a stale duration on the activity row can't corrupt the schedule.
        if is_milestone:
            duration_hrs = 0
        progress_pct = max(0, min(100, float(activity.get('physical_progress_pct') or 0)))
        remaining_days = activity.get('remaining_duration_days')
        remaining_hrs = int(float(remaining_days) * hours_per_day) if remaining_days is not None else duration_hrs
        actual_start = _xer_datetime(activity.get('actual_start'), XER_WORK_START)
        actual_finish = _xer_datetime(activity.get('actual_finish'), XER_WORK_START if is_milestone else XER_WORK_END)
        task_status = 'TK_Complete' if actual_finish else ('TK_Active' if actual_start else 'TK_NotStart')
        float_days = activity.get('total_float_days') or 0
        float_hrs = int(float_days * hours_per_day)

        start_raw = activity.get('start_date')
        finish_raw = activity.get('finish_date')
        early_start = _xer_datetime(start_raw, XER_WORK_START)
        early_end = _xer_datetime(finish_raw, XER_WORK_START if is_milestone else XER_WORK_END)

        late_start = early_start
        late_end = early_end
        if float_days:
            shifted_start = _xer_add_working_days(start_raw, int(float_days), XER_WORK_START)
            shifted_end = _xer_add_working_days(finish_raw, int(float_days), XER_WORK_START if is_milestone else XER_WORK_END)
            if shifted_start and shifted_end:
                late_start, late_end = shifted_start, shifted_end

        lines.append(_xer_row(
            task_id_map[activity['id']], XER_PROJ_ID, wbs_id, XER_CLNDR_ID, progress_pct, 'N', 1, 'N', 'N',
            'CP_Drtn',
            'TT_Mile' if is_milestone else 'TT_Task',
            'DT_FixedDUR2', task_status,
            str(activity['id'])[:XER_TASK_CODE_MAX_LEN],
            (activity.get('name') or '')[:XER_TASK_NAME_MAX_LEN],
            '',                                # rsrc_id (none — resources managed downstream)
            float_hrs, float_hrs, remaining_hrs,
            0, 0, 0, duration_hrs, 0, 0, 0,    # work / equip quantities (planning-only export)
            '',                                # cstr_date
            actual_start, actual_finish,
            late_start, late_end, '',
            early_start, early_end,
            early_start, early_end,            # restart_date / reend_date
            early_start, early_end,            # target_start_date / target_end_date
            late_start, late_end,              # rem_late_start / rem_late_end
            '', 'PT_Normal',                   # cstr_type / priority_type
            '', '', '', '',                    # suspend/resume/float_path/float_path_order
            task_guid_map[activity['id']], '', '', '',
            '', 0, 0, '', '',
            now_stamp, now_stamp, export_user, export_user, '', '',
        ))

    # ── TASKPRED ──
    # Collect and de-duplicate predecessor links from both sources so the
    # imported schedule always carries a fully-connected logic network —
    # P6's scheduler otherwise flags every unlinked activity as "open"
    # and won't compute a critical path.
    lines.append('%T\tTASKPRED')
    lines.append('%F\ttask_pred_id\ttask_id\tpred_task_id\tproj_id\tpred_proj_id\tpred_type\t'
                  'lag_hr_cnt\tcomments\tfloat_path\taref\tarls')

    seen_pred_keys = set()
    pred_rows = []

    def _record_pred(succ_id, pred_id, rel_type, lag_hrs=0):
        succ = task_id_map.get(succ_id)
        pred = task_id_map.get(pred_id)
        if not succ or not pred or succ == pred:
            return
        key = (succ, pred, rel_type)
        if key in seen_pred_keys:
            return
        seen_pred_keys.add(key)
        pred_rows.append((succ, pred, rel_type, lag_hrs))

    for link in logic_matrix or []:
        rel = _XER_PRED_TYPE.get((link.get('type') or 'FS').upper(), 'PR_FS')
        _record_pred(
            link.get('activity_id') or link.get('task_id'),
            link.get('predecessor_id') or link.get('pred_task_id'),
            rel,
            int((link.get('lag_days') or 0) * hours_per_day),
        )
    for activity in activities:
        for pred in activity.get('predecessors') or []:
            rel = _XER_PRED_TYPE.get((pred.get('type') or 'FS').upper(), 'PR_FS')
            _record_pred(
                activity['id'], pred.get('id'), rel,
                int((pred.get('lag_days') or 0) * hours_per_day),
            )

    for seq, (succ, pred, rel_type, lag_hrs) in enumerate(pred_rows, start=1):
        lines.append(_xer_row(
            seq, succ, pred, XER_PROJ_ID, XER_PROJ_ID, rel_type, lag_hrs,
            '', '', '', '',
        ))

    lines.append('%E')

    text = XER_LINE_SEP.join(lines) + XER_LINE_SEP
    try:
        return text.encode('cp1252')
    except UnicodeEncodeError:
        return text.encode('utf-8')
