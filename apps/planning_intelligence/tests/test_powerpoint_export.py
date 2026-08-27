from io import BytesIO
from types import SimpleNamespace
from unittest import TestCase
from unittest.mock import patch

from pptx import Presentation

from apps.planning_intelligence.services.export_utils import generation_to_pptx_bytes


class PowerPointExportTests(TestCase):
    def _generation(self, *, with_milestone=True):
        project = SimpleNamespace(
            name='RADAI Demo Project',
            phase='FEED',
            client='Demo Client',
            location='Abu Dhabi',
            effective_date='2026-08-27',
            duration_months=6,
        )
        milestone = {
            'id': 'M-001',
            'name': 'Project kickoff',
            'finish_date': '2026-09-01',
            'is_critical': True,
            'is_milestone': True,
        }
        return SimpleNamespace(
            project=project,
            version=1,
            activities=[milestone] if with_milestone else [],
            wbs=[{'level': 1, 'code': '1', 'name': 'Engineering'}],
            eddr=[{
                'discipline': 'Process',
                'deliverable_name': 'Design Basis',
                'final_issue_date': '2026-09-30',
            }],
            manhours={
                'by_discipline': [{
                    'discipline_name': 'Process',
                    'responsible_role': 'Engineer',
                    'man_days': 12,
                    'man_hours': 96,
                }],
                'grand_total_man_hours': 96,
            },
            validation=[{'severity': 'pass', 'message': 'Schedule is valid'}],
            intelligence={'disciplines': {'process': {}}},
            milestones=[milestone] if with_milestone else [],
            narrative='The schedule establishes a controlled delivery sequence.',
        )

    @patch(
        'apps.planning_intelligence.services.export_utils.PPTX_TEMPLATE_PATH',
        'missing-corporate-template.pptx',
    )
    def test_generates_valid_widescreen_deck_without_template_asset(self):
        content = generation_to_pptx_bytes(self._generation())

        self.assertTrue(content.startswith(b'PK'))
        presentation = Presentation(BytesIO(content))
        self.assertEqual(len(presentation.slides), 11)
        self.assertEqual(presentation.slides[0].shapes.title.text, 'RADAI Demo Project')
        self.assertEqual(presentation.slides[-1].shapes.title.text, 'Thank You')
        self.assertGreater(len(content), 20_000)

    @patch(
        'apps.planning_intelligence.services.export_utils.PPTX_TEMPLATE_PATH',
        'missing-corporate-template.pptx',
    )
    def test_omits_optional_milestone_slide_when_no_milestones_exist(self):
        content = generation_to_pptx_bytes(self._generation(with_milestone=False))
        presentation = Presentation(BytesIO(content))
        titles = [slide.shapes.title.text for slide in presentation.slides if slide.shapes.title]

        self.assertEqual(len(presentation.slides), 10)
        self.assertNotIn('Key Milestones', titles)
        self.assertIn('Executive Summary', titles)

    @patch(
        'apps.planning_intelligence.services.export_utils.PPTX_TEMPLATE_PATH',
        'missing-corporate-template.pptx',
    )
    def test_paginates_long_executive_summary_without_losing_text(self):
        generation = self._generation()
        generation.narrative = '\n'.join(
            f'Paragraph {number}: ' + ('Detailed planning information ' * 18)
            for number in range(1, 9)
        ) + '\nFINAL-NARRATIVE-MARKER'

        content = generation_to_pptx_bytes(generation)
        presentation = Presentation(BytesIO(content))
        summary_slides = [
            slide for slide in presentation.slides
            if slide.shapes.title and slide.shapes.title.text.startswith('Executive Summary')
        ]
        summary_text = '\n'.join(
            shape.text
            for slide in summary_slides
            for shape in slide.shapes
            if getattr(shape, 'has_text_frame', False)
        )

        self.assertGreater(len(summary_slides), 1)
        self.assertEqual(summary_slides[0].shapes.title.text, f'Executive Summary (1/{len(summary_slides)})')
        self.assertIn('FINAL-NARRATIVE-MARKER', summary_text)
        for slide in summary_slides:
            body_text = '\n'.join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, 'has_text_frame', False) and shape != slide.shapes.title
            )
            self.assertLessEqual(len(body_text), 800)
